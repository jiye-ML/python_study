# 使用Python制作PPT
from pptx import Presentation

# 实例化Presentation
ppt = Presentation()
# 使用ppt自带的模板, ppt自带了1-48中模板
title_slide_layout = ppt.slide_layouts[0]
# 新建一页幻灯片
slide = ppt.slides.add_slide(title_slide_layout)
# 设置标题
title = slide.shapes.title
# 使用placeholdes索引获取一页幻灯片中的元素
subtitle = slide.placeholders[1]
# 设置text
title.text = "Hello, World!"
subtitle.text = "PPT自动化 from python-pptx!"
# 保存文件
ppt.save('demo.pptx')


# 新建一页幻灯片
title_slide_layout = ppt.slide_layouts[5]
slide = ppt.slides.add_slide(title_slide_layout)

# 添加新文本框
from pptx.util import Inches

# 预设位置及大小
left, top = Inches(0.5), Inches(1.6)
width, height = Inches(4), Inches(1)
# left，top为相对位置，width，height为文本框大小
textbox = slide.shapes.add_textbox(left, top, width, height)
# 文本框中文字
textbox.text = 'this is a new textbox'
# 在新文本框中添加段落
new_para = textbox.text_frame.add_paragraph()
# 段落文字
new_para.text = 'this is second paragraph in textbox'

# 文件路径
img_path ='pic1.jpg'
# 预设位置及大小
left, top, width, height = Inches(0.5), Inches(2.5), Inches(8), Inches(5)
# 在指定位置按预设值添加图片
pic = slide.shapes.add_picture(img_path, left, top, width, height)

ppt.save('demo.pptx')


