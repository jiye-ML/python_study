# 根据每页的摘要和图片，生成PPT
from pptx import Presentation
# 添加新文本框
from pptx.util import Inches, Pt
import os


# 找配图
def find_pic(filename):
	filenames = []
	#print(filename)
	files= os.listdir(pic_path) #得到文件夹下的所有文件名称
	for file in files: #遍历文件夹
		if file.startswith(filename):
			filenames.append(file)
			#print(file)
	return filenames

#文件夹目录
title_text = '不同场景的不同拍摄技巧'
#title_text = "Python入门"
pic_path = './images'
text_path = "./text"
# 实例化Presentation
ppt = Presentation()

def add_slide(ppt, text, pics):
	# 使用ppt自带的模板, ppt自带了1-48中模板
	bullet_slide_layout = ppt.slide_layouts[3]
	slide = ppt.slides.add_slide(bullet_slide_layout)
	shapes = slide.shapes
	# 设置标题
	title = slide.shapes.title
	title.text = title_text

	# 使用placeholdes索引获取一页幻灯片中的元素
	tf = slide.placeholders[1]
	tf = tf.text_frame
	p = tf.add_paragraph()
	p.text = text
	p.font.size = Pt(18)
	# 右侧放图片
	if len(pics) > 0:
		# 文件路径
		img_path = pic_path + '/' + pics[0]
		# 预设位置及大小
		#left, top, width, height = Inches(0.5), Inches(2.5), Inches(8), Inches(5)
		# 放置到右侧
		left, top, width = Pt(370), Pt(170), Pt(320)
		# 在指定位置按预设值添加图片
		#pic = slide.shapes.add_picture(img_path, left, top, width, height)
		pic = slide.shapes.add_picture(img_path, left, top, width)

files= os.listdir(text_path) #得到文件夹下的所有文件名称
#s = []
for file in files: #遍历文件夹
    if not os.path.isdir(file): #判断是否是文件夹，不是文件夹才打开
        f = open(text_path+"/"+file) #打开文件
        text = f.read()
        # 找配图
        pics = find_pic(file[:-4])
        # 创建slide
        add_slide(ppt, text, pics)

ppt.save('book.pptx')
