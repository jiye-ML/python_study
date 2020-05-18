# 打开Demo.pptx 添加
from pptx import Presentation

# 实例化Presentation
prs = Presentation()
# 使用ppt自带的模板, ppt自带了1-48中模板
bullet_slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(bullet_slide_layout)
shapes = slide.shapes

# 标题
title_shape = shapes.title
title_shape.text = 'Office Automation by Python'

# 正文
body_shape = shapes.placeholders[1]
# 获取文本框
tf = body_shape.text_frame
tf.text = 'Excel Operation by Python'

# 给文本框添加0级段落
p = tf.add_paragraph()
p.text = 'PPT Operation by Python'
p.level = 0
# 添加1级段落
p = tf.add_paragraph()
p.text = 'Use pptx tools to control PPT files'
p.level = 1
# 添加2级段落
p = tf.add_paragraph()
p.text = 'add_paragraph() function usage'
p.level = 2

prs.save('pptx_usage.pptx')